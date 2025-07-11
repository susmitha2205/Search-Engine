import os
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
import PyPDF2
import docx
from pptx import Presentation
import json
import streamlit as st
import ollama
from io import BytesIO
from typing import List, Dict, Tuple, Optional

# Initialize model and globals
model = SentenceTransformer('sentence-transformers/msmarco-bert-base-dot-v5')
dimension = 768

class DocumentReader:
    @staticmethod
    def read_pdf(file) -> str:
        """Read text from PDF file."""
        try:
            if hasattr(file, 'read'):
                file.seek(0)
                reader = PyPDF2.PdfReader(BytesIO(file.read()))
            else:
                reader = PyPDF2.PdfReader(file)
            
            text = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
            return ' '.join(text)
        except Exception as e:
            st.error(f"Error reading PDF {getattr(file, 'name', 'file')}: {str(e)}")
            return ""

    @staticmethod
    def read_docx(file) -> str:
        """Read text from DOCX file."""
        try:
            if hasattr(file, 'read'):
                file.seek(0)
                doc = docx.Document(BytesIO(file.read()))
            else:
                doc = docx.Document(file)
            return ' '.join([para.text for para in doc.paragraphs if para.text])
        except Exception as e:
            st.error(f"Error reading DOCX {getattr(file, 'name', 'file')}: {str(e)}")
            return ""

    @staticmethod
    def read_pptx(file) -> str:
        """Read text from PPTX file."""
        try:
            if hasattr(file, 'read'):
                file.seek(0)
                prs = Presentation(BytesIO(file.read()))
            else:
                prs = Presentation(file)
            
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return ' '.join(text)
        except Exception as e:
            st.error(f"Error reading PPTX {getattr(file, 'name', 'file')}: {str(e)}")
            return ""

    @staticmethod
    def read_txt(file) -> str:
        """Read text from TXT file."""
        try:
            if hasattr(file, 'read'):
                file.seek(0)
                return file.read().decode('utf-8')
            else:
                with open(file, 'r', encoding='utf-8') as f:
                    return f.read()
        except Exception as e:
            st.error(f"Error reading TXT {getattr(file, 'name', 'file')}: {str(e)}")
            return ""

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """Split text into overlapping chunks."""
    words = text.split()
    if len(words) <= chunk_size:
        return [' '.join(words)]
    
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = ' '.join(words[i:i + chunk_size])
        chunks.append(chunk)
    return chunks

def index_documents(files: List) -> Tuple[Optional[faiss.Index], Optional[List[Dict]]]:
    """Create a FAISS index from document contents."""
    index = faiss.IndexFlatIP(dimension)
    metadata = []
    documents = []
    reader = DocumentReader()

    for file in files:
        name = file.name
        content = ""
        
        if name.lower().endswith('.pdf'):
            content = reader.read_pdf(file)
        elif name.lower().endswith('.docx'):
            content = reader.read_docx(file)
        elif name.lower().endswith('.pptx'):
            content = reader.read_pptx(file)
        elif name.lower().endswith('.txt'):
            content = reader.read_txt(file)
        else:
            st.warning(f"Unsupported file type: {name}")
            continue

        if not content.strip():
            st.warning(f"No content extracted from {name}")
            continue

        chunks = chunk_text(content)
        for i, chunk in enumerate(chunks):
            documents.append(chunk)
            metadata.append({'file': name, 'chunk': i, 'content': chunk})

    if not documents:
        st.error("No valid content to index.")
        return None, None

    embeddings = model.encode(documents, show_progress_bar=True)
    embeddings = np.array(embeddings, dtype='float32')
    index.add(embeddings)
    return index, metadata

def semantic_search(query: str, index: faiss.Index, metadata: List[Dict], k: int = 5) -> List[Dict]:
    """Perform semantic search using the FAISS index."""
    query_embedding = model.encode([query])[0].astype('float32')
    distances, indices = index.search(np.array([query_embedding]), k)
    
    results = []
    for score, idx in zip(distances[0], indices[0]):
        if idx < len(metadata):
            results.append({
                'score': float(score),
                'file': metadata[idx]['file'],
                'chunk': metadata[idx]['chunk'],
                'content': metadata[idx]['content']
            })
    return results

def generate_answer(query: str, context: List[str], model_name: str = 'gemma:2b') -> str:
    """Generate an answer using the provided context. Returns either:
    - LLM-generated answer
    - Summary of relevant passages if generation fails"""
    if not context:
        return "No relevant context found to answer the question."
    
    joined_context = "\n\n".join([f"Source {i+1}:\n{ctx}" for i, ctx in enumerate(context)])
    
    prompt = f"""You are a helpful AI assistant. Answer the question based ONLY on the following context. 
If the answer isn't in the context, say you don't know. Be concise and accurate.

Context:
{joined_context}

Question: {query}

Answer:"""
    
    try:
        response = ollama.generate(
            model=model_name,
            prompt=prompt,
            options={'temperature': 0.1}
        )
        return response['response'].strip()
    except Exception as e:
        # Create summary from references when generation fails
        summary = "\n".join([f"• {ctx[:300]}..." for ctx in context])
        return summary

def main():
    st.set_page_config(page_title="Generative Search Engine for Local Files", layout="wide", page_icon="📚")
    st.title("📚 Generative Search Engine for Local Files")
    st.markdown("Upload your documents and ask questions about their content.")
    
    # Initialize session state
    if 'indexed' not in st.session_state:
        st.session_state.indexed = False
    if 'files' not in st.session_state:
        st.session_state.files = []
    
    with st.sidebar:
        st.header("Settings")
        chunk_size = st.slider("Chunk Size (words)", 200, 1000, 500)
        overlap = st.slider("Chunk Overlap (words)", 0, 200, 50)
        top_k = st.slider("Number of chunks to retrieve", 1, 10, 3)
        llm_model = st.selectbox("LLM Model", ['gemma:2b', 'mistral', 'llama2'])
    
    st.header("1. Upload Documents")
    uploaded_files = st.file_uploader(
        "Choose files (PDF, DOCX, PPTX, TXT)",
        type=['pdf', 'docx', 'pptx', 'txt'],
        accept_multiple_files=True
    )
    
    if st.button("Index Documents") and uploaded_files:
        with st.spinner("Processing documents..."):
            index, metadata = index_documents(uploaded_files)
            if index and metadata:
                st.session_state.index = index
                st.session_state.metadata = metadata
                st.session_state.files = uploaded_files
                st.session_state.indexed = True
                st.success(f"Indexed {len(uploaded_files)} documents with {len(metadata)} chunks!")
            else:
                st.error("Failed to index documents.")

    if st.session_state.get('indexed', False):
        st.header("2. Ask a Question")
        question = st.text_area("Enter your question about the documents:", height=100)
        
        if st.button("Get Answer") and question.strip():
            with st.spinner("Searching for relevant information..."):
                results = semantic_search(
                    question,
                    st.session_state.index,
                    st.session_state.metadata,
                    k=top_k
                )
                
                if results:
                    context = [res['content'] for res in results]
                    answer = generate_answer(question, context, llm_model)
                    
                    st.subheader("Answer")
                    st.write(answer)
                    
                    st.subheader("References")
                    for i, res in enumerate(results, 1):
                        with st.expander(f"Source {i}: {res['file']} (Score: {res['score']:.2f})"):
                            st.write(res['content'])
                else:
                    st.warning("No relevant information found in the documents.")

if __name__ == '__main__':
    main()

